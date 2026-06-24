import os
import re
import logging
import pdfplumber
import pandas as pd
from pptx import Presentation
from openpyxl import load_workbook
from openpyxl.styles import PatternFill

# 1. Setup Logging for cleaner execution tracking
logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")

# 2. Global Configuration Constants
# To be run locally, so I have not encry
EXCEL_FILE = "DocumentDataTracker.xlsx"
TARGET_SHEET = "MasterTracker"
INPUT_FOLDER = "./test_files"

# Definitive 9-Column Architecture + Confidence Scoring
COLUMNS = [ # Fill columns here based on master sheet on xlsx file
    "File Name", "Project Name", "Objective", 
    "Field 1", "Field 2", "Field 3", "Field 4","Field 5", 
    "Status", "Confidence Score", "Review Status"
]

def init_excel():
    """Verifies your pre-existing Excel file is present before running."""
    if not os.path.exists(EXCEL_FILE):
        # We raise an error now instead of making a file, matching your workflow requirement
        raise FileNotFoundError(
            f"❌ Critical Error: The baseline tracker file '{EXCEL_FILE}' was not found. "
            f"Please ensure your pre-existing Excel file is placed in this directory."
        )
    else:
        logging.info(f"✔ Pre-existing tracker connected successfully: {EXCEL_FILE}")

def append_to_excel(data_dict):
    """
    Cross-validates column headers to map data to the correct visual layout 
    positions dynamically, then merges data into existing or new records.
    """
    book = load_workbook(EXCEL_FILE)
    sheet = book[TARGET_SHEET] if TARGET_SHEET in book.sheetnames else book.active

    # 1. DYNAMIC HEADER CROSS-VALIDATION
    # Create a map of { "Column Header Name": column_index_integer }
    header_map = {}
    for col_idx in range(1, sheet.max_column + 1):
        header_val = sheet.cell(row=1, column=col_idx).value
        if header_val:
            header_map[str(header_val).strip()] = col_idx

    # Verify all expected columns are present in the target sheet
    missing_columns = [col for col in COLUMNS if col not in header_map]
    if missing_columns:
        raise KeyError(
            f"❌ Header Validation Failed! The following required column headers "
            f"are missing from your Excel sheet: {missing_columns}. "
            f"Please verify your column naming matches exactly."
        )

    # 2. MAP FILE NAME TO FIND EXISTING ROW (Using validated column index)
    file_col_idx = header_map["File Name"]
    existing_records = {}
    for row_idx in range(2, sheet.max_row + 1):
        cell_value = sheet.cell(row=row_idx, column=file_col_idx).value
        if cell_value:
            existing_records[str(cell_value).strip()] = row_idx

    target_filename = str(data_dict["File Name"]).strip()
    
    if target_filename in existing_records:
        target_row = existing_records[target_filename]
        logging.info(f"🔄 Cross-Validated: Merging data into existing row {target_row} for '{target_filename}'.")
    else:
        target_row = sheet.max_row + 1
        logging.info(f"➕ Cross-Validated: Appending fresh entry at row {target_row} for '{target_filename}'.")

    # 3. INTELLIGENT FIELD POPULATION (Writes to mapped headers only)
    for col_name in COLUMNS:
        target_col_idx = header_map[col_name]
        current_cell = sheet.cell(row=target_row, column=target_col_idx)
        new_value = data_dict.get(col_name)
        
        # Merge Rule: Only populate if current cell is empty/NaN/UNKNOWN
        if current_cell.value is None or str(current_cell.value).strip() in ["", "UNKNOWN", "NaN"]:
            current_cell.value = new_value
            
            # Apply HITL formatting alerts to newly filled blocks
            if data_dict["Review Status"] == "Pending Review":
                numeric_score = int(float(data_dict["Confidence Score"].replace("%", "")))
                color = "FCE4D6" if numeric_score < 60 else "FFF2CC"
                current_cell.fill = PatternFill(start_color=color, end_color=color, fill_type="solid")
            
    book.save(EXCEL_FILE)
    
def extract_text_from_pdf(file_path):
    """
    Extracts raw text strings from PDF layers.
    Uses layout=True to preserve visual whitespace structural positioning,
    preventing text from jumping across columns or wrapping incorrectly.
    """
    full_text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                # ENHANCED: layout=True forces the engine to respect multi-column bounding areas
                text = page.extract_text(layout=True)
                if text:
                    full_text += text + "\n"
    except Exception as e:
        logging.error(f"Failed to read non-linear PDF {file_path}: {e}")
    return full_text

def extract_text_from_pptx(file_path):
    """Extracts text elements natively from all floating shapes inside a PowerPoint file."""
    full_text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    full_text += shape.text + "\n"
    except Exception as e:
        logging.error(f"Failed to read PowerPoint layout {file_path}: {e}")
    return full_text

def calculate_confidence(extracted_data):
    """Calculates a localized data certainty score out of 100%."""
    score = 0
    critical_fields = ["Project Name", "Objective", "Status"]
    
    # 💡 UPDATED: Added "Field 5" to the optional fields list
    optional_fields = ["Field 1", "Field 2", "Field 3", "Field 4", "Field 5"]
    
    # Tier 1: Field Presence Weights (Max 60 points)
    for field in critical_fields:
        if extracted_data[field] != "UNKNOWN":
            score += 15  # 45 points total for critical targets
            
    for field in optional_fields:
        if extracted_data[field] != "UNKNOWN":
            score += 3.0  # 💡 UPDATED: 3 points each across 5 fields (15 total)
            
    # Tier 2: Length Anomaly Validation (Max 20 points)
    length_penalty = False
    if len(extracted_data["Project Name"]) > 100 or len(extracted_data["Status"]) > 50:
        length_penalty = True
        
    if not length_penalty and extracted_data["Project Name"] != "UNKNOWN":
        score += 20
        
    # Tier 3: Data Integrity Sanitization Check (Max 20 points)
    if not any(char in extracted_data["Project Name"] for char in ["{", "}", "[", "]", "<", ">"]):
        score += 20
        
    return f"{min(score, 100)}%"

def parse_text_with_hitl(text, filename):
    """Processes unformatted multi-column layout strings via Regex pattern maps."""
    # Anchors target terms. If text wraps weirdly due to formatting, regex handles line isolation.
   # 💡 UPDATED: Added anchor pattern for Field 5
    patterns = {
        "Project Name": r"Project:\s*(.*)",
        "Objective": r"Objective:\s*(.*)",
        "Field 1": r"Field 1:\s*(.*)",
        "Field 2": r"Field 2:\s*(.*)",
        "Field 3": r"Field 3:\s*(.*)",
        "Field 4": r"Field 4:\s*(.*)",
        "Field 5": r"Field 5:\s*(.*)",  # <-- ADD THIS LINE
        "Status": r"Status:\s*(.*)"
    }
    
    extracted = {"File Name": filename}
    
    for key, regex in patterns.items():
        match = re.search(regex, text, re.IGNORECASE)
        extracted[key] = match.group(1).strip() if match else "UNKNOWN"
        
    # Inject Confidence Calculations
    extracted["Confidence Score"] = calculate_confidence(extracted)
    numeric_score = int(float(extracted["Confidence Score"].replace("%", "")))
    
    # Human-in-the-Loop Routing Rule
    # Triggers manual pipeline validation if metrics score drops below 85% or core keys match UNKNOWN
    if numeric_score < 85 or extracted["Project Name"] == "UNKNOWN" or extracted["Status"] == "UNKNOWN":
        extracted["Review Status"] = "Pending Review"
    else:
        extracted["Review Status"] = "Approved"
        
    return extracted

def run_pipeline():
    """Execution wrapper running pipeline steps iteratively across files."""
    init_excel()
    
    if not os.path.exists(INPUT_FOLDER) or not os.listdir(INPUT_FOLDER):
        logging.warning(f"Drop your PDFs/PPTXs into the '{INPUT_FOLDER}' directory to begin data runs.")
        os.makedirs(INPUT_FOLDER, exist_ok=True)
        return

    for file in os.listdir(INPUT_FOLDER):
        file_path = os.path.join(INPUT_FOLDER, file)
        raw_text = ""
        
        if file.endswith(".pdf"):
            raw_text = extract_text_from_pdf(file_path)
        elif file.endswith(".pptx"):
            raw_text = extract_text_from_pptx(file_path)
        else:
            continue  # Skips unmapped extension schemas
            
        data = parse_text_with_hitl(raw_text, file)
        append_to_excel(data)
        logging.info(f"Processed: {file} | Score: {data['Confidence Score']} | Flag Status: {data['Review Status']}")

if __name__ == "__main__":
    run_pipeline()