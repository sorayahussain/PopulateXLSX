import os
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

INPUT_FOLDER = "./test_files"
os.makedirs(INPUT_FOLDER, exist_ok=True)

def make_pdf(filename, text_lines):
    """Generates a real local PDF file with custom text lines."""
    path = os.path.join(INPUT_FOLDER, filename)
    c = canvas.Canvas(path, pagesize=letter)
    c.setFont("Helvetica", 10)
    
    y = 750  # Start near the top of the page
    for line in text_lines:
        c.drawString(50, y, line)
        y -= 20  # Move down for the next line
        
    c.save()
    print(f"Generated PDF: {path}")

def make_pptx(filename, slide_texts):
    """Generates a real local PPTX slide deck with custom text blocks."""
    path = os.path.join(INPUT_FOLDER, filename)
    prs = Presentation()
    
    for text_content in slide_texts:
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank slide layout
        # Add a text frame box visually on the slide
        txBox = slide.shapes.add_textbox(100000, 100000, 5000000, 4000000)
        tf = txBox.text_frame
        tf.text = text_content
        
    prs.save(path)
    print(f"Generated PPTX: {path}")

# --- Data Profiles ---
clean_pdf_text = [
    "================================================================================",
    "ANNUAL PERFORMANCE SUMMARY REPORT",
    "================================================================================",
    "Project: Project Alpha Evolution",
    "Objective: To modernize the internal legacy data warehousing infrastructure across departments.",
    "Status: Active",
    "--------------------------------------------------------------------------------",
    "Field 1: 45 Days Remaining",
    "Field 2: $120,000 Budget Allocated",
    "Field 3: Tier 1 Engineering Team Assigned",
    "Field 4: 87% Completion Rate",
]

# Simulating layout confusion by printing columns side-by-side with spacing
chaotic_pdf_text = [
    "================================================================================",
    "EXECUTIVE BRIEF                                QUARTERLY INITIATIVES",
    "================================================================================",
    "Project: Project Beta Phoenix                  Field 1: Phase 3 Deployment",
    "Status: Pending Approval                       Field 2: Global Region Operations",
    "--------------------------------------------------------------------------------",
    "Objective: Transitioning consumer face         Field 3: Cloud Infrastructure Core",
    "applications to a serverless architecture      Field 4: Enterprise Compliance Verified",
    "================================================================================",
]

broken_ppt_slides = [
    "INFOGRAPHIC DASHBOARD GRAPHICS (VISUAL ONLY)\n\nObjective: Expand international user registration pipelines by Q4.",
    "DATA BLOCKS:\nField 1: Region NA & EU\nField 2: Localization Complete\nField 3: 12.5M Active Nodes\nField 4: API Connection Stable"
]

# Run Generator
if __name__ == "__main__":
    make_pdf("clean_linear_report.pdf", clean_pdf_text)
    make_pdf("chaotic_multicolumn_report.pdf", chaotic_pdf_text)
    make_pptx("broken_infographic_mock.pptx", broken_ppt_slides)
    print("\n🎉 All test files successfully placed inside './test_files/'!")